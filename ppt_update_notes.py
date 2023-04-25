# Open powerpoint file

from pptx import Presentation
import collections
import collections.abc

file_path = r'/home/andy/git-python/extract_standard_info/extract_table/slide_test.pptx'
powerpoint = Presentation(file_path)

# Loop through all slides and print the title
for slide in powerpoint.slides:
    print(slide.shapes[0].text)



